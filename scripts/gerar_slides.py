"""
gerar_slides.py
===============
Gera apresentação .pptx do TP com layouts visuais (cards, colunas, destaques).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Paleta — azul UFMG / amarelo institucional / cinza profissional
PRIMARIA   = RGBColor(0x1F, 0x3A, 0x5F)   # azul escuro
SECUNDARIA = RGBColor(0xE8, 0xB4, 0x0F)   # amarelo UFMG
ACENTO     = RGBColor(0x2E, 0x86, 0xAB)   # azul médio
SUCESSO    = RGBColor(0x2E, 0x7D, 0x32)   # verde
ALERTA     = RGBColor(0xC0, 0x50, 0x4D)   # vermelho suave
ROXO       = RGBColor(0x6A, 0x1B, 0x9A)   # roxo
TEXTO      = RGBColor(0x21, 0x25, 0x29)
TEXTO_LEVE = RGBColor(0x55, 0x5B, 0x66)
FUNDO_CARD = RGBColor(0xF5, 0xF7, 0xFA)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)

LARGURA = Inches(13.333)
ALTURA  = Inches(7.5)


# ---------------------------------------------------------------------------
# HELPERS DE LAYOUT
# ---------------------------------------------------------------------------

def _set_text(tf, texto, *, size=20, bold=False, italic=False,
              color=TEXTO, align=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font_name
    return r


def add_titulo_pagina(slide, texto, *, subtitulo=None):
    # barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  LARGURA, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARIA
    bar.line.fill.background()

    # acento amarelo
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.85),
                                  LARGURA, Inches(0.06))
    acc.fill.solid()
    acc.fill.fore_color.rgb = SECUNDARIA
    acc.line.fill.background()

    # título
    txt = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.7))
    _set_text(txt.text_frame, texto, size=28, bold=True, color=BRANCO)

    if subtitulo:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5))
        _set_text(sub.text_frame, subtitulo, size=15, italic=True, color=TEXTO_LEVE)
        return Inches(1.6)
    return Inches(1.2)


def add_card(slide, *, left, top, width, height,
             titulo, texto, num=None,
             cor_borda=PRIMARIA, cor_titulo=PRIMARIA, fundo=FUNDO_CARD,
             tamanho_titulo=16, tamanho_texto=13):
    """Card colorido com título e texto. Opcionalmente exibe número em círculo."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fundo
    box.line.color.rgb = cor_borda
    box.line.width = Pt(1.5)
    box.adjustments[0] = 0.06

    pad = Inches(0.2)
    inner_left = left + pad
    inner_width = width - 2 * pad

    if num is not None:
        d = Inches(0.5)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       left + Inches(0.15), top + Inches(0.15),
                                       d, d)
        circ.fill.solid()
        circ.fill.fore_color.rgb = cor_borda
        circ.line.fill.background()
        circ_tf = circ.text_frame
        circ_tf.margin_left = Inches(0)
        circ_tf.margin_right = Inches(0)
        circ_tf.margin_top = Inches(0.02)
        circ_tf.margin_bottom = Inches(0)
        p = circ_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(num)
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = BRANCO
        title_left = inner_left + Inches(0.5)
        title_width = inner_width - Inches(0.5)
    else:
        title_left = inner_left
        title_width = inner_width

    tbox = slide.shapes.add_textbox(title_left, top + Inches(0.15),
                                     title_width, Inches(0.45))
    _set_text(tbox.text_frame, titulo, size=tamanho_titulo, bold=True, color=cor_titulo)

    body = slide.shapes.add_textbox(inner_left, top + Inches(0.7),
                                     inner_width, height - Inches(0.85))
    btf = body.text_frame
    btf.word_wrap = True
    if isinstance(texto, list):
        # filtra strings vazias para evitar bullets fantasmas
        items_validos = [it for it in texto if it.strip()]
        for i, item in enumerate(items_validos):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(4)
            r = p.add_run()
            r.text = "• " + item
            r.font.size = Pt(tamanho_texto)
            r.font.color.rgb = TEXTO
    else:
        _set_text(btf, texto, size=tamanho_texto, color=TEXTO)


def add_metric(slide, *, left, top, width, height, numero, rotulo, cor=PRIMARIA):
    """Card destaque para um número grande com rótulo embaixo."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = cor
    box.line.fill.background()
    box.adjustments[0] = 0.06

    n = slide.shapes.add_textbox(left, top + Inches(0.15),
                                  width, Emu(int(height * 0.55)))
    _set_text(n.text_frame, str(numero), size=40, bold=True,
              color=BRANCO, align=PP_ALIGN.CENTER)

    r_box = slide.shapes.add_textbox(left, top + Emu(int(height * 0.6)),
                                      width, Emu(int(height * 0.35)))
    _set_text(r_box.text_frame, rotulo, size=12, color=BRANCO,
              align=PP_ALIGN.CENTER)


def add_imagem_centralizada(slide, path, *, top=1.7, height=5.4):
    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top), height=Inches(height))
    pic.left = int((LARGURA - pic.width) / 2)


def add_footer(slide, n_atual, num_total):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    _set_text(box.text_frame,
              f"TP Banco de Dados · DCC/UFMG · {n_atual}/{num_total}",
              size=9, italic=True, color=PRIMARIA, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# BUILD
# ---------------------------------------------------------------------------

def build(out: Path):
    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA
    blank = prs.slide_layouts[6]

    # ============ Slide 1: CAPA ============
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), LARGURA, ALTURA)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARIA; bg.line.fill.background()
    deco = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.5), Inches(0),
                               Inches(4), Inches(7.5))
    deco.fill.solid(); deco.fill.fore_color.rgb = ACENTO; deco.line.fill.background()
    deco.rotation = 180
    faixa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.35),
                                Inches(10), Inches(0.06))
    faixa.fill.solid(); faixa.fill.fore_color.rgb = SECUNDARIA; faixa.line.fill.background()

    eb = s.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(10), Inches(0.4))
    _set_text(eb.text_frame, "TRABALHO PRÁTICO · BANCO DE DADOS · DCC/UFMG",
              size=12, bold=True, color=SECUNDARIA)

    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11), Inches(1.6))
    _set_text(tb.text_frame, "Do Mundo Real ao Insight",
              size=56, bold=True, color=BRANCO)

    sb = s.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(11), Inches(0.8))
    _set_text(sb.text_frame, "Modelagem, implementação e análise do CPGF",
              size=22, italic=True, color=SECUNDARIA)

    autores = s.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(11), Inches(1.7))
    tf2 = autores.text_frame
    nomes = ["Lucas Soares Benfica", "Vinícius Cardoso Antunes", "Pedro Soares Pinto"]
    for i, nome in enumerate(nomes):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(2)
        rr = p.add_run(); rr.text = nome
        rr.font.size = Pt(20); rr.font.color.rgb = BRANCO

    rod = s.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(11), Inches(0.4))
    _set_text(rod.text_frame,
              "Prof. Pedro H. Barros · 2026",
              size=13, italic=True, color=SECUNDARIA)

    # ============ Slide 2: AGENDA (cards numerados 2x3) ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "Agenda",
                              subtitulo="O caminho dos próximos ~14 minutos")

    cards = [
        (1, "Domínio escolhido",     "Por que o CPGF é um caso ideal",     ACENTO),
        (2, "Modelo Conceitual",     "ER/EER · 7 entidades · restrições",  PRIMARIA),
        (3, "Modelo Relacional",     "Mapeamento + Normalização BCNF",     ROXO),
        (4, "Implementação SQL",     "DDL + Carga em PostgreSQL",          SUCESSO),
        (5, "Consultas & Insights",  "5 perguntas analíticas + bônus",     ALERTA),
        (6, "Conclusão",             "Aprendizados e próximos passos",     PRIMARIA),
    ]
    col_w = Inches(4.0)
    row_h = Inches(2.45)
    for i, (n, titulo, desc, cor) in enumerate(cards):
        col = i % 3
        row = i // 3
        left = Inches(0.5) + col * (col_w + Inches(0.2))
        top = top_c + row * (row_h + Inches(0.2))
        add_card(s, left=left, top=top, width=col_w, height=row_h,
                 titulo=titulo, texto=desc, num=n,
                 cor_borda=cor, cor_titulo=cor,
                 tamanho_titulo=16, tamanho_texto=13)

    # ============ Slide 3: DOMÍNIO (2 colunas) ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "Por que o CPGF?",
                              subtitulo="Cartão de Pagamento do Governo Federal · Portal da Transparência")

    add_card(s, left=Inches(0.5), top=top_c,
             width=Inches(5.8), height=Inches(5.0),
             titulo="O que é",
             texto=[
                "Cartão de crédito corporativo de UGs federais",
                "Transação = órgão, portador, favorecido, data, tipo, valor",
                "Publicado mensalmente pelo Portal da Transparência (CGU)",
                "Dado pessoal anonimizado conforme LGPD",
             ],
             cor_borda=PRIMARIA, cor_titulo=PRIMARIA,
             tamanho_titulo=18, tamanho_texto=14)

    razoes = [
        ("Volume",     "50-100 mil transações/mês",            ACENTO),
        ("Estrutura",  "4+ entidades (hierarquia + EER)",      PRIMARIA),
        ("Abertura",   "Acesso público via LAI",               SUCESSO),
        ("Relevância", "Controle social do gasto público",     ALERTA),
    ]
    col_w = Inches(3.1)
    row_h = Inches(2.4)
    for i, (tit, desc, cor) in enumerate(razoes):
        col = i % 2
        row = i // 2
        left = Inches(6.6) + col * (col_w + Inches(0.15))
        top = top_c + row * (row_h + Inches(0.2))
        add_card(s, left=left, top=top, width=col_w, height=row_h,
                 titulo=tit, texto=desc,
                 cor_borda=cor, cor_titulo=cor,
                 tamanho_titulo=16, tamanho_texto=13)

    # ============ Slide 4: ER (imagem) ============
    s = prs.slides.add_slide(blank)
    add_titulo_pagina(s, "Etapa 1 — Modelo Conceitual (ER/EER)",
                      subtitulo="7 entidades · notação ER clássica + extensão EER (especialização disjunta total)")
    add_imagem_centralizada(s, "docs/diagrama_er.png", top=1.7, height=5.4)

    # ============ Slide 5: RESTRIÇÕES (3 grupos temáticos) ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "Restrições de Integridade",
                              subtitulo="Nove restrições do domínio · embarcadas no schema PostgreSQL")

    grupos = [
        ("Valores", ALERTA, [
            "R1  valor ≠ 0",
            "R2  valor > 0 (negativo só em reversões)",
        ]),
        ("Datas & Formato", ACENTO, [
            "R3  data próxima do mês-ano do extrato",
            "R4  mês ∈ {1..12}; ano ≥ 2003",
            "R5  CPF do portador na máscara ***.NNN.NNN-**",
        ]),
        ("Identificação", SUCESSO, [
            "R6  11 dígitos = PF; 14 = PJ; outros = NI",
            "R7  integridade referencial (4 FKs em TRANSACAO)",
            "R8  hierarquia UG → Subordinado → Superior",
            "R9  discriminador EER ∈ {PF, PJ, NI}",
        ]),
    ]
    col_w = Inches(4.05)
    for i, (tit, cor, items) in enumerate(grupos):
        left = Inches(0.5) + i * (col_w + Inches(0.2))
        add_card(s, left=left, top=top_c, width=col_w, height=Inches(5.0),
                 titulo=tit, texto=items,
                 cor_borda=cor, cor_titulo=cor,
                 tamanho_titulo=18, tamanho_texto=14)

    # ============ Slide 6: RELACIONAL (imagem) ============
    s = prs.slides.add_slide(blank)
    add_titulo_pagina(s, "Etapa 2 — Esquema Relacional",
                      subtitulo="9 relações com PK/FK explícitas em PostgreSQL")
    add_imagem_centralizada(s, "docs/diagrama_relacional.png", top=1.7, height=5.4)

    # ============ Slide 7: NORMALIZAÇÃO (fluxo) ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "Normalização até BCNF",
                              subtitulo="Por que decompor o CSV bruto")

    fluxo = [
        ("CSV bruto", ALERTA, [
            "Tabela plana de 15 colunas",
            "DFs transitivas internas",
            "Viola 3FN",
        ]),
        ("Anomalias geradas", PRIMARIA, [
            "Update: renomear ministério = milhões de updates",
            "Insert: nova UG depende de transação",
            "Redundância massiva",
        ]),
        ("9 relações em BCNF", SUCESSO, [
            "Cada DF transitiva → tabela própria",
            "JOIN natural reconstrói o CSV",
            "Sem dependência não-chave → não-chave",
        ]),
    ]
    col_w = Inches(3.8)
    h = Inches(4.3)
    seta_y = top_c + Emu(int(h / 2)) - Inches(0.25)
    for i, (tit, cor, items) in enumerate(fluxo):
        left = Inches(0.5) + i * (col_w + Inches(0.4))
        add_card(s, left=left, top=top_c, width=col_w, height=h,
                 titulo=tit, texto=items,
                 cor_borda=cor, cor_titulo=cor,
                 tamanho_titulo=17, tamanho_texto=13)
        if i < 2:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      left + col_w + Inches(0.05), seta_y,
                                      Inches(0.3), Inches(0.5))
            arr.fill.solid(); arr.fill.fore_color.rgb = TEXTO_LEVE
            arr.line.fill.background()

    final = s.shapes.add_textbox(Inches(0.5), top_c + h + Inches(0.3),
                                  Inches(12.3), Inches(0.5))
    _set_text(final.text_frame,
              "Decisão de NÃO desnormalizar: consistência > microsegundos no contexto de LAI",
              size=14, italic=True, color=TEXTO_LEVE, align=PP_ALIGN.CENTER)

    # ============ Slide 8 ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "Etapa 3 — Implementação em PostgreSQL",
                              subtitulo="DDL + pipeline ETL + limpeza · totalmente reprodutível")

    metricas = [
        ("9",     "tabelas",             PRIMARIA),
        ("9",     "restrições R1-R9",    ACENTO),
        ("6",     "índices auxiliares",  ROXO),
        ("9.898", "linhas carregadas",   SUCESSO),
        ("101",   "descartes auditados", ALERTA),
    ]
    m_w = Inches(2.35)
    m_h = Inches(1.3)
    gap = Inches(0.12)
    total_w = len(metricas) * m_w + (len(metricas) - 1) * gap
    base_left = Emu(int((LARGURA - total_w) / 2))
    for i, (num, lbl, cor) in enumerate(metricas):
        left = base_left + i * (m_w + gap)
        add_metric(s, left=left, top=top_c, width=m_w, height=m_h,
                    numero=num, rotulo=lbl, cor=cor)

    card_top = top_c + m_h + Inches(0.4)
    card_h = Inches(3.2)
    add_card(s, left=Inches(0.5), top=card_top, width=Inches(6.1), height=card_h,
             titulo="sql/01_ddl.sql",
             texto=[
                "CREATE TABLE das 9 relações",
                "PKs, FKs com ON UPDATE CASCADE / ON DELETE RESTRICT",
                "CHECK constraints embarcadas (R1, R3, R4...)",
                "Índices em data, UG, portador, favorecido, tipo, ano-mês",
             ],
             cor_borda=PRIMARIA, cor_titulo=PRIMARIA,
             tamanho_titulo=17, tamanho_texto=13)
    add_card(s, left=Inches(6.8), top=card_top, width=Inches(6.0), height=card_h,
             titulo="scripts/carregar_dados.py",
             texto=[
                "Lê CSV em LATIN1, separador ;",
                "Tipa, valida, desduplica, descarta linhas problemáticas",
                "Popula 9 tabelas em ordem topológica",
                "Idempotente · suporta PostgreSQL e SQLite",
             ],
             cor_borda=SUCESSO, cor_titulo=SUCESSO,
             tamanho_titulo=17, tamanho_texto=13)

    # ============ Slides 9-13 (graficos) ============
    for q, titulo, subt in [
        ("q1_gasto_por_ministerio", "P1 — Quais ministérios mais gastam via CPGF?",
              "JOIN de 3 tabelas + GROUP BY + agregações"),
        ("q2_serie_temporal",       "P2 — Há sazonalidade nos gastos?",
              "GROUP BY composto + agregação temporal"),
        ("q3_top_favorecidos",      "P3 — Top-10 favorecidos e concentração",
              "JOIN + GROUP BY + HAVING + LIMIT"),
        ("q4_tipo_transacao",       "P4 — Distribuição por tipo de transação",
              "Saques em espécie merecem monitoramento"),
        ("q5_top_portadores",       "P5 — Outliers por ministério",
              "CTE + RANK() OVER (PARTITION BY ministerio ORDER BY gasto DESC)"),
    ]:
        s = prs.slides.add_slide(blank)
        add_titulo_pagina(s, titulo, subtitulo=subt)
        add_imagem_centralizada(s, f"docs/graficos/{q}.png", top=1.7, height=5.4)

    # ============ Slide 14: BONUS ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "Bonus — Além do mínimo exigido",
                              subtitulo="Detecção de anomalias e conexão com a LAI")

    add_card(s, left=Inches(0.5), top=top_c,
             width=Inches(6.1), height=Inches(5.0),
             titulo="P6 — Detecção de anomalias",
             texto=[
                "Critério: |valor| > média + 3·desvio-padrão",
                "< 1% das linhas mas parcela desproporcional do valor",
                "Triagem automatizada — útil ao controle interno",
                "Atende: 'técnicas além do escopo'",
             ],
             cor_borda=ALERTA, cor_titulo=ALERTA,
             tamanho_titulo=18, tamanho_texto=14)

    add_card(s, left=Inches(6.8), top=top_c,
             width=Inches(6.0), height=Inches(5.0),
             titulo="Conexão com a LAI",
             texto=[
                "Lei 12.527/2011 — Acesso à Informação",
                "CPGF: dado anonimizado, público, para controle social",
                "Cruzar dimensões = trabalho da CGU",
                "Atende: 'cidadania de dados'",
             ],
             cor_borda=ACENTO, cor_titulo=ACENTO,
             tamanho_titulo=18, tamanho_texto=14)

    # ============ Slide 15: CONCLUSAO ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "O que entregamos",
                              subtitulo="Síntese dos resultados e próximos passos")

    metricas = [
        ("9", "relações em BCNF",     PRIMARIA),
        ("5", "consultas analíticas", ACENTO),
        ("1", "window function",      ROXO),
        ("6", "visualizações",        SUCESSO),
        ("3", "bônus contemplados",   ALERTA),
    ]
    m_w = Inches(2.35)
    m_h = Inches(1.4)
    gap = Inches(0.12)
    total_w = len(metricas) * m_w + (len(metricas) - 1) * gap
    base_left = Emu(int((LARGURA - total_w) / 2))
    for i, (num, lbl, cor) in enumerate(metricas):
        left = base_left + i * (m_w + gap)
        add_metric(s, left=left, top=top_c, width=m_w, height=m_h,
                    numero=num, rotulo=lbl, cor=cor)

    card_top = top_c + m_h + Inches(0.45)
    add_card(s, left=Inches(0.5), top=card_top,
             width=Inches(6.1), height=Inches(2.9),
             titulo="Principal aprendizado",
             texto=[
                "A qualidade do modelo conceitual",
                "determina a clareza das consultas.",
                "Um modelo mal normalizado teria forçado",
                "workarounds em cada query analítica.",
             ],
             cor_borda=PRIMARIA, cor_titulo=PRIMARIA,
             tamanho_titulo=18, tamanho_texto=14)
    add_card(s, left=Inches(6.8), top=card_top,
             width=Inches(6.0), height=Inches(2.9),
             titulo="Próximos passos",
             texto=[
                "Cruzamento com a base CEIS (Sanções)",
                "Carga incremental mensal automatizada",
                "Dashboard interativo alinhado à LAI",
                "Detecção de padrões via técnicas estatísticas",
             ],
             cor_borda=SUCESSO, cor_titulo=SUCESSO,
             tamanho_titulo=18, tamanho_texto=14)

    # ============ Slide 16: OBRIGADO ============
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), LARGURA, ALTURA)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARIA; bg.line.fill.background()
    deco = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(5.5),
                               Inches(13.333), Inches(2.0))
    deco.fill.solid(); deco.fill.fore_color.rgb = ACENTO; deco.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2.0))
    _set_text(tb.text_frame, "Obrigado!", size=84, bold=True,
              color=BRANCO, align=PP_ALIGN.CENTER)
    sb = s.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12), Inches(0.8))
    _set_text(sb.text_frame, "Perguntas?", size=32, italic=True,
              color=SECUNDARIA, align=PP_ALIGN.CENTER)

    # ============ Slide 17: IA ============
    s = prs.slides.add_slide(blank)
    top_c = add_titulo_pagina(s, "Declaração de uso de IA",
                              subtitulo="Conforme permitido pelo enunciado")

    add_card(s, left=Inches(0.5), top=top_c,
             width=Inches(6.1), height=Inches(4.6),
             titulo="✓  Onde a IA foi usada",
             texto=[
                "Apoio na redação do relatório e destes slides",
                "Sugestões de boas práticas em SQL e modelagem ER/EER",
                "Geração de código boilerplate para o pipeline ETL",
                "Revisão ortográfica",
             ],
             cor_borda=SUCESSO, cor_titulo=SUCESSO,
             tamanho_titulo=18, tamanho_texto=14)

    add_card(s, left=Inches(6.8), top=top_c,
             width=Inches(6.0), height=Inches(4.6),
             titulo="✗  O que NÃO foi delegado",
             texto=[
                "Decisões de modelagem e normalização",
                "Interpretação dos resultados e discussão crítica",
                "Escolha do dataset e perguntas analíticas",
                "Validação dos achados",
             ],
             cor_borda=ALERTA, cor_titulo=ALERTA,
             tamanho_titulo=18, tamanho_texto=14)

    rod = s.shapes.add_textbox(Inches(0.5), top_c + Inches(4.85),
                                Inches(12.3), Inches(0.5))
    _set_text(rod.text_frame,
              "Ferramenta: assistente de IA generativa. Bibliotecas no README: pandas, matplotlib, graphviz, python-docx, python-pptx",
              size=11, italic=True, color=TEXTO_LEVE, align=PP_ALIGN.CENTER)

    # ----------------------------------------------------------------------
    # Footers
    # ----------------------------------------------------------------------
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx in (1, 16):
            continue
        add_footer(slide, idx, total)

    prs.save(str(out))
    print(f"OK - apresentacao gerada em {out} ({total} slides)")


if __name__ == "__main__":
    out = Path("slides/Apresentacao_TP_BD.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    build(out)
